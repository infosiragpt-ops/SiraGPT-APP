"""Real validator fixture tests. No validator or external-tool mocks.

These run directly on the developer host and do NOT attest container isolation.
"""
import hashlib
import importlib.util
import io
import json
import os
from pathlib import Path
import stat
import tempfile
import unittest
import zipfile

SCRIPT = Path(__file__).parents[1] / 'src/modules/doc-sandbox/validation/validator.py'
spec = importlib.util.spec_from_file_location('validator', SCRIPT)
validator = importlib.util.module_from_spec(spec)
spec.loader.exec_module(validator)


def archive(parts, compress=zipfile.ZIP_STORED):
    output = io.BytesIO()
    with zipfile.ZipFile(output, 'w', compression=compress) as result:
        for name, content in parts:
            result.writestr(name, content)
    return output.getvalue()


def docx(text='Original', middle=' bold', tail=' tail'):
    return archive([
        ('[Content_Types].xml', '''<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>'''),
        ('_rels/.rels', '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>'''),
        ('word/document.xml', f'''<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>{text}</w:t></w:r><w:r><w:rPr><w:b/></w:rPr><w:t xml:space="preserve">{middle}</w:t></w:r><w:r><w:t xml:space="preserve">{tail}</w:t></w:r></w:p><w:p><w:r><w:t>Untouched paragraph</w:t></w:r></w:p><w:sectPr><w:pgSz w:w="11906" w:h="16838"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/></w:sectPr></w:body></w:document>'''),
    ])


def xlsx(value='10', recalc=False):
    return archive([
        ('[Content_Types].xml', '''<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/><Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/></Types>'''),
        ('_rels/.rels', '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>'''),
        ('xl/workbook.xml', f'''<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="Budget" sheetId="1" r:id="rId1"/></sheets><calcPr fullCalcOnLoad="{'1' if recalc else '0'}"/></workbook>'''),
        ('xl/_rels/workbook.xml.rels', '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>'''),
        ('xl/worksheets/sheet1.xml', f'''<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheetData><row r="1"><c r="A1"><v>{value}</v></c><c r="B1"><f>A1*2</f><v>20</v></c></row></sheetData></worksheet>'''),
    ])


class FixtureTest(unittest.TestCase):
    def setUp(self):
        self.directory = tempfile.TemporaryDirectory(prefix='doc-validator-test-')
        self.root = Path(self.directory.name)

    def tearDown(self):
        self.directory.cleanup()

    def write(self, name, data):
        path = self.root / name
        path.write_bytes(data)
        return path

    def failure(self, code, action):
        with self.assertRaises(validator.ValidationFailure) as raised:
            action()
        self.assertEqual(raised.exception.code, code)

    def plan(self, original, inventory, replacement):
        unit = inventory['units'][0]
        return {'schemaVersion': 1, 'mode': 'preserve', 'outputName': inventory['name'],
                'inputHashes': {'one': inventory['sha256']}, 'notPossible': [],
                'edits': [] if replacement is None else [
                    {'id': 'edit-1', 'kind': 'text', 'inputId': 'one', 'part': unit['part'],
                     'locator': unit['locator'], 'before': unit['text'], 'after': replacement}]}

    def request(self, original, output, name, plan):
        return {'command': 'validate', 'inputs': [{'id': 'one', 'path': str(original), 'name': name}],
                'outputPath': str(output), 'plan': plan, 'artifactDir': str(self.root / 'artifacts')}


class ZipSecurity(FixtureTest):
    def test_zip_slip(self):
        for entry in ('../bad.xml', '/bad.xml', 'C:/bad.xml', '..\\bad.xml', '%2e%2e/bad.xml', 'a/./bad.xml'):
            with self.subTest(entry=entry):
                self.failure('ZIP_PATH_UNSAFE', lambda: validator.safe_zip(archive([(entry, '<a/>')])) )

    def test_duplicate(self):
        self.failure('ZIP_DUPLICATE', lambda: validator.safe_zip(archive([('a.xml', '<a/>'), ('a.xml', '<b/>')])))

    def test_symbolic_link(self):
        entry = zipfile.ZipInfo('link')
        entry.create_system = 3
        entry.external_attr = (stat.S_IFLNK | 0o777) << 16
        self.failure('ZIP_SPECIAL_FILE', lambda: validator.safe_zip(archive([(entry, '/etc/passwd')])))

    def test_zip_bomb(self):
        self.failure('ZIP_BOMB', lambda: validator.safe_zip(archive([('bomb.xml', '0' * 1000000)], zipfile.ZIP_DEFLATED)))

    def test_corrupt_zip(self):
        self.failure('INVALID_ZIP', lambda: validator.safe_zip(b'not zip'))

    def test_crc_checked(self):
        data = archive([('a.xml', 'HELLO')]).replace(b'HELLO', b'HXLLO')
        self.failure('INVALID_ZIP', lambda: validator.safe_zip(data))

    def test_xxe(self):
        for content in (b'<!DOCTYPE x [<!ENTITY y SYSTEM "file:///etc/passwd">]><x>&y;</x>',
                        '<!DOCTYPE x [<!ENTITY y "secret">]><x>&y;</x>'.encode('utf-16')):
            self.failure('XML_DTD_FORBIDDEN', lambda: validator.xml(content))

    def test_malformed_xml(self):
        self.failure('INVALID_XML', lambda: validator.xml(b'<x>'))

    def test_fake_docx(self):
        file = self.write('fake.docx', b'%PDF-1.7\n')
        self.failure('MIME_MISMATCH', lambda: validator.inspect(file, file.name))

    def test_mime_is_not_just_pk(self):
        file = self.write('fake.docx', archive([('anything', 'data')]))
        self.failure('OOXML_MISSING_PART', lambda: validator.inspect(file, file.name))

    def test_missing_relationship(self):
        data = docx().replace(b'Target="word/document.xml"', b'Target="word/nonexist.xml"')
        # CRC remains invalid when editing ZIP bytes, so create a proper malicious archive.
        parts, _ = validator.safe_zip(docx())
        parts['_rels/.rels'] = parts['_rels/.rels'].replace(b'word/document.xml', b'word/missing.xml')
        file = self.write('broken.docx', archive(parts.items()))
        self.failure('INVALID_RELATIONSHIP', lambda: validator.inspect(file, file.name))

    def test_external_content(self):
        parts, _ = validator.safe_zip(docx())
        parts['_rels/.rels'] = parts['_rels/.rels'].replace(b'Target="word/document.xml"', b'Target="https://example.com/image" TargetMode="External"')
        file = self.write('external.docx', archive(parts.items()))
        self.failure('EXTERNAL_CONTENT', lambda: validator.inspect(file, file.name))

    def test_hidden_macro(self):
        parts, _ = validator.safe_zip(docx())
        parts['word/vbaProject.bin'] = b'macro'
        parts['[Content_Types].xml'] = parts['[Content_Types].xml'].replace(b'</Types>', b'<Default Extension="bin" ContentType="application/octet-stream"/></Types>')
        file = self.write('macro.docx', archive(parts.items()))
        self.failure('ACTIVE_CONTENT_UNSUPPORTED', lambda: validator.inspect(file, file.name))

    def test_recipe_inspection_parses_but_never_executes(self):
        marker = self.root / 'must-not-exist'
        data = archive([('recipe/01_edit.py', f'from pathlib import Path\nPath({str(marker)!r}).write_text("NO")\n')])
        file = self.write('recipe.zip', data)
        result = validator.run({'command': 'inspect_recipe', 'inputs': [{'id': 'r', 'path': str(file), 'name': 'recipe.zip'}]})
        self.assertEqual(result['recipe']['scripts'], ['recipe/01_edit.py'])
        self.assertFalse(marker.exists())

    def test_recipe_path_traversal_rejected(self):
        file = self.write('recipe.zip', archive([('../evil.py', 'pass')]))
        self.failure('ZIP_PATH_UNSAFE', lambda: validator.run({'command': 'inspect_recipe',
            'inputs': [{'id': 'r', 'path': str(file), 'name': 'recipe.zip'}]}))

    def test_recipe_nested_binary_rejected(self):
        file = self.write('recipe.zip', archive([('nested.zip', b'PK\x03\x04'), ('edit.py', 'pass')]))
        self.failure('RECIPE_FILE_UNSUPPORTED', lambda: validator.run({'command': 'inspect_recipe',
            'inputs': [{'id': 'r', 'path': str(file), 'name': 'recipe.zip'}]}))

    def test_recipe_python_syntax_rejected(self):
        file = self.write('recipe.zip', archive([('edit.py', 'def [invalid:')]))
        self.failure('RECIPE_SYNTAX', lambda: validator.run({'command': 'inspect_recipe',
            'inputs': [{'id': 'r', 'path': str(file), 'name': 'recipe.zip'}]}))


class ExactEdits(FixtureTest):
    def documents(self, data=None):
        original = self.write('before.docx', docx())
        output = self.write('after.docx', data or docx('Edited'))
        before = validator.inspect(original, 'document.docx')
        after = validator.inspect(output, 'document.docx')
        op = {'kind': 'replace_text', 'part': 'word/document.xml', 'locator': before['units'][0]['locator'],
              'before': 'Original', 'after': 'Edited'}
        return original, output, before, after, [op]

    def test_exact_leaf_retains_mixed_styles(self):
        original, output, before, after, edits = self.documents()
        self.assertEqual(validator.structural(original, output, before, after, edits)['exactNodeChanges'], 1)
        self.assertTrue(validator.textual(before, after, edits)['exact'])

    def test_style_change_inside_authorized_paragraph_is_rejected(self):
        parts, _ = validator.safe_zip(docx('Edited'))
        parts['word/document.xml'] = parts['word/document.xml'].replace(b'<w:b/>', b'<w:i/>')
        args = self.documents(archive(parts.items()))
        self.failure('UNAUTHORIZED_XML_CHANGE', lambda: validator.structural(*args))

    def test_other_text_in_same_part_is_rejected(self):
        args = self.documents(docx('Edited', middle='hacked'))
        self.failure('UNAUTHORIZED_XML_CHANGE', lambda: validator.structural(*args))

    def test_unplanned_attribute(self):
        parts, _ = validator.safe_zip(docx('Edited'))
        parts['word/document.xml'] = parts['word/document.xml'].replace(b'w:top="1440"', b'w:top="1500"')
        args = self.documents(archive(parts.items()))
        self.failure('UNAUTHORIZED_XML_CHANGE', lambda: validator.structural(*args))

    def test_noop_is_byte_exact(self):
        original, output, before, after, edits = self.documents(docx())
        self.assertEqual(validator.structural(original, output, before, after, [])['exactNodeChanges'], 0)
        self.failure('NOOP_CHANGED', lambda: validator.structural(*self.documents()[:4], []))

    def test_empty_plan_rejects_zip_only_changes(self):
        original, output, before, after, _ = self.documents(archive(validator.safe_zip(docx())[0].items(), zipfile.ZIP_DEFLATED))
        self.failure('NOOP_CHANGED', lambda: validator.structural(original, output, before, after, []))

    def test_partial_substring_before_does_not_authorize_unit(self):
        original, output, before, after, edits = self.documents()
        edits[0]['before'] = 'Orig'
        self.failure('PLAN_TEXT_MISMATCH', lambda: validator.check_plan(before, {'inputSha256': before['sha256'], 'mode': 'preserve', 'operations': edits}))

    def test_xpath_injection_not_evaluated(self):
        original, output, before, after, edits = self.documents()
        edits[0]['locator'] = '//w:t'
        self.failure('PLAN_LOCATOR', lambda: validator.check_plan(before, {'inputSha256': before['sha256'], 'mode': 'preserve', 'operations': edits}))

    def test_visual_context_is_derived_from_all_runs(self):
        original, output, before, after, edits = self.documents()
        self.assertEqual(validator.visual_operations(original, before, edits), [{'before': 'Original bold tail', 'after': 'Edited bold tail'}])


class PlainValidation(FixtureTest):
    def test_real_pipeline_text_formats(self):
        cases = [('txt', 'Before\r\n', 'After\r\n'), ('md', '# Before\n', '# After\n'),
                 ('json', '{"year": 2026}\n', '{"year": 2027}\n'),
                 ('csv', 'name,year\nA,2026\n', 'name,year\nA,2027\n'),
                 ('html', '<html><p>Before</p></html>\n', '<html><p>After</p></html>\n')]
        for fmt, first, second in cases:
            with self.subTest(format=fmt):
                original = self.write('before.' + fmt, first.encode())
                output = self.write('after.' + fmt, second.encode())
                inventory = validator.inspect(original, 'document.' + fmt)
                report = validator.run(self.request(original, output, inventory['name'], self.plan(original, inventory, second)))['report']
                self.assertTrue(report['passed'], report)
                self.assertEqual(len(report['levels']), 4)
                self.assertFalse(report['levels'][1]['applicable'])
                self.assertFalse(report['levels'][1]['passed'])
                self.assertTrue((self.root / 'artifacts/text-diff.json').is_file())

    def test_invalid_json_is_not_a_valid_edited_result(self):
        original = self.write('a.json', b'{"year":2026}')
        output = self.write('b.json', b'{"year":2027')
        inventory = validator.inspect(original, 'a.json')
        report = validator.run(self.request(original, output, 'a.json', self.plan(original, inventory, '{"year":2027')))['report']
        self.assertFalse(report['passed'])
        self.assertEqual(report['levels'][0]['details']['code'], 'TEXT_STRUCTURE')

    def test_unplanned_extra_text_fails_exact_diff(self):
        original = self.write('a.txt', b'before')
        output = self.write('b.txt', b'after EXTRA')
        inventory = validator.inspect(original, 'a.txt')
        report = validator.run(self.request(original, output, 'a.txt', self.plan(original, inventory, 'after')))['report']
        self.assertFalse(report['passed'])
        self.assertEqual(report['levels'][-1]['details']['code'], 'TEXT_DIFF_UNPLANNED')

    def test_preserve_encoding(self):
        original = self.write('a.txt', 'before'.encode('utf-16'))
        output = self.write('b.txt', 'after'.encode('utf-8'))
        inventory = validator.inspect(original, 'a.txt')
        report = validator.run(self.request(original, output, 'a.txt', self.plan(original, inventory, '\ufeffafter')))['report']
        self.assertFalse(report['passed'])
        self.assertEqual(report['levels'][0]['details']['code'], 'ENCODING_CHANGED')

    def test_preserve_crlf(self):
        original = self.write('a.txt', b'before\r\n')
        output = self.write('b.txt', b'after\n')
        inventory = validator.inspect(original, 'a.txt')
        report = validator.run(self.request(original, output, 'a.txt', self.plan(original, inventory, 'after\n')))['report']
        self.assertEqual(report['levels'][0]['details']['code'], 'LINE_ENDINGS_CHANGED')

    def test_noop_text_pipeline(self):
        original = self.write('a.txt', b'unchanged\n')
        inventory = validator.inspect(original, 'a.txt')
        report = validator.run(self.request(original, original, 'a.txt', self.plan(original, inventory, None)))['report']
        self.assertTrue(report['passed'], report)

    def test_active_html_rejected(self):
        file = self.write('a.html', b'<html><script>alert(1)</script></html>')
        self.failure('HTML_ACTIVE_CONTENT', lambda: validator.inspect(file, 'a.html'))


class ToolsAndRealOpening(FixtureTest):
    def test_missing_tool_does_not_pass(self):
        self.failure('TOOL_UNAVAILABLE', lambda: validator.command(['siragpt-required-nonexistent-validator']))

    def test_real_libreoffice_and_poppler_noop(self):
        original = self.write('a.docx', docx())
        inventory = validator.inspect(original, 'a.docx')
        report = validator.run(self.request(original, original, 'a.docx', self.plan(original, inventory, None)))['report']
        self.assertTrue(report['passed'], report)
        self.assertTrue(all(item['passed'] and item['applicable'] for item in report['levels']))

    def test_real_libreoffice_and_poppler_edit(self):
        original = self.write('a.docx', docx())
        output = self.write('b.docx', docx('Changed'))
        inventory = validator.inspect(original, 'a.docx')
        report = validator.run(self.request(original, output, 'a.docx', self.plan(original, inventory, 'Changed')))['report']
        self.assertTrue(report['passed'], report)

    def test_ambiguous_visual_locator_fails_closed(self):
        self.failure('VISUAL_LOCATOR_AMBIGUOUS', lambda: validator.phrase_regions(
            [[{'text': 'word', 'box': (1, 1, 5, 5)}, {'text': 'word', 'box': (5, 1, 9, 5)}]], 'word'))

    def test_pdf_expected_merge_rotate_and_overlay_uses_real_libraries(self):
        from reportlab.pdfgen import canvas
        from pypdf import PdfReader
        inputs = []
        for number in range(2):
            file = self.root / f'{number}.pdf'
            drawing = canvas.Canvas(str(file), pagesize=(400, 500), invariant=1)
            drawing.drawString(40, 450, f'Original {number}')
            drawing.save()
            inputs.append({'id': str(number), 'path': str(file), 'name': file.name})
        edits = [{'kind': 'pdf_merge', 'inputIds': ['0', '1']},
                 {'kind': 'pdf_overlay', 'inputId': '0', 'page': 1, 'text': 'Page 1', 'x': 40, 'y': 20, 'fontSize': 10},
                 {'kind': 'pdf_rotate', 'inputId': '1', 'pages': [1], 'degrees': 90}]
        output = validator.expected_pdf(inputs, edits, self.root / 'expected.pdf')
        reader = PdfReader(output)
        self.assertEqual(len(reader.pages), 2)
        self.assertIn('Original 0', reader.pages[0].extract_text())
        self.assertIn('Page 1', reader.pages[0].extract_text())
        self.assertEqual(reader.pages[1].rotation, 90)
        self.assertTrue(validator.pdf_structure(output, output)['formsAndAnnotationsChecked'])

    def test_workbook_recalc_and_derived_visual_baseline(self):
        original = self.write('a.xlsx', xlsx())
        output = self.write('b.xlsx', xlsx('30', recalc=True))
        inventory = validator.inspect(original, 'a.xlsx')
        plan = self.plan(original, inventory, '30')
        plan['edits'][0]['kind'] = 'cell'
        report = validator.run(self.request(original, output, 'a.xlsx', plan))['report']
        self.assertTrue(report['passed'], report)
        self.assertTrue(report['levels'][1]['details']['formulasRecalculated'])
        self.assertEqual(report['levels'][2]['details']['baseline'], 'independently-applied-cell-edits-and-recalculation')

    def test_workbook_cannot_leave_stale_formula_cache_without_recalc(self):
        original = self.write('a.xlsx', xlsx())
        output = self.write('b.xlsx', xlsx('30'))
        inventory = validator.inspect(original, 'a.xlsx')
        report = validator.run(self.request(original, output, 'a.xlsx', self.plan(original, inventory, '30')))['report']
        self.assertEqual(report['levels'][0]['details']['code'], 'FORMULA_RECALC_REQUIRED')

    def test_real_presentation_opening_and_text_edit(self):
        # python-pptx is fixture generation only, never the document editor.
        from pptx import Presentation
        from pptx.util import Inches
        original = self.root / 'a.pptx'
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
        box.text = 'Professional original'
        slide.notes_slide.notes_text_frame.text = 'Speaker original'
        deck.save(original)
        inventory = validator.inspect(original, 'a.pptx')
        unit = next(u for u in inventory['units'] if u['part'] == 'ppt/slides/slide1.xml' and u['text'] == 'Professional original')
        note = next(u for u in inventory['units'] if u['part'] == 'ppt/notesSlides/notesSlide1.xml' and u['text'] == 'Speaker original')
        parts, order = validator.safe_zip(original.read_bytes())
        root = validator.xml(parts[unit['part']])
        validator.locate(root, unit['locator']).text = 'Professional revised'
        parts[unit['part']] = validator.etree.tostring(root)
        root = validator.xml(parts[note['part']])
        validator.locate(root, note['locator']).text = 'Speaker revised'
        parts[note['part']] = validator.etree.tostring(root)
        output = self.write('b.pptx', archive((name, parts[name]) for name in order))
        plan = {'schemaVersion': 1, 'mode': 'preserve', 'outputName': 'a.pptx', 'inputHashes': {'one': inventory['sha256']},
                'notPossible': [], 'edits': [{'id': 'e1', 'kind': 'text', 'inputId': 'one', 'part': unit['part'],
                    'locator': unit['locator'], 'before': unit['text'], 'after': 'Professional revised'},
                    {'id': 'e2', 'kind': 'text', 'inputId': 'one', 'part': note['part'],
                     'locator': note['locator'], 'before': note['text'], 'after': 'Speaker revised'}]}
        report = validator.run(self.request(original, output, 'a.pptx', plan))['report']
        self.assertTrue(report['passed'], report)
        self.assertTrue(report['levels'][2]['details']['notes']['applicable'])

    def test_real_visual_rejects_extra_content_outside_plan(self):
        from reportlab.pdfgen import canvas
        for name, extra in [('before.pdf', False), ('after.pdf', True)]:
            drawing = canvas.Canvas(str(self.root / name), pagesize=(400, 500), invariant=1)
            drawing.drawString(40, 450, 'Unchanged document')
            if extra:
                drawing.rect(20, 20, 150, 100, fill=1)
            drawing.save()
        work = self.root / 'visual'
        work.mkdir()
        self.failure('VISUAL_UNAUTHORIZED_CHANGE', lambda: validator.visual(
            self.root / 'before.pdf', self.root / 'after.pdf', [], work))


if __name__ == '__main__':
    unittest.main(verbosity=2)
