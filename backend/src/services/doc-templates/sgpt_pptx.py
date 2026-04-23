"""sgpt_pptx — professional deck templates for siraGPT.

Opinionated deck builder:
  · Widescreen 16:9 at 13.33 × 7.5 in.
  · Single cohesive palette (navy + cream + terracotta accent) —
    every slide uses the same colours.
  · Title slide, agenda, section dividers, content (1-col / 2-col /
    big-stat), quote, closing.
  · Footer with page number + project name on every slide.
  · Speaker notes slot exposed per slide.

Usage:
    from sgpt_pptx import Deck
    d = Deck(title='Defensa de Tesis', subtitle='Autor · 2025', palette='tesis_upn')
    d.cover()
    d.agenda(['Introducción', 'Metodología', 'Resultados', 'Conclusiones'])
    d.section('1. Introducción')
    d.text_slide('Problema', ['Bullet 1', 'Bullet 2'])
    d.big_stat('68%', 'reducción de tiempos con SMED')
    d.two_column('Comparación', left_title='Antes', left=['a','b'], right_title='Después', right=['c','d'])
    d.thanks()
    d.save(OUT_PATH)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy


# ─── Palettes ────────────────────────────────────────────────────────────

PALETTES = {
    'tesis_upn': {
        'bg':        'FFFFFF',
        'surface':   'F7F3EA',
        'primary':   '1F3A68',   # navy
        'accent':    'C05621',   # terracotta
        'text':      '1A1918',
        'muted':     '6B7280',
    },
    'defense': {
        'bg':        'FFFFFF',
        'surface':   'F9FAFB',
        'primary':   '0F172A',
        'accent':    '6366F1',
        'text':      '111827',
        'muted':     '6B7280',
    },
    'pitch': {
        'bg':        '0F172A',
        'surface':   '1E293B',
        'primary':   'F97316',
        'accent':    'FBBF24',
        'text':      'F1F5F9',
        'muted':     '94A3B8',
    },
}


def _rgb(hexcode):
    return RGBColor.from_string(hexcode.upper().lstrip('#'))


# ─── Deck ────────────────────────────────────────────────────────────────

class Deck:
    def __init__(self, *, title, subtitle=None, author=None, institution=None,
                 date=None, palette='tesis_upn'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.palette = PALETTES.get(palette, PALETTES['tesis_upn'])
        self.title = title
        self.subtitle = subtitle
        self.author = author
        self.institution = institution
        self.date = date
        self._page_num = 0

    # ─── Helpers ─────────────────────────────────────────────────────────

    def _new_slide(self, with_footer=True):
        blank = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank)
        # background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height,
        )
        bg_shape.line.fill.background()
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _rgb(self.palette['bg'])
        # send to back
        spTree = bg_shape._element.getparent()
        spTree.remove(bg_shape._element)
        spTree.insert(2, bg_shape._element)
        if with_footer:
            self._page_num += 1
            self._add_footer(slide, self._page_num)
        return slide

    def _add_footer(self, slide, page_num):
        # Thin accent line + page number + project title on bottom edge.
        y = self.prs.slide_height - Inches(0.55)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), y + Inches(0.05),
            Inches(2.0), Inches(0.03),
        )
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(self.palette['accent'])

        tb_title = slide.shapes.add_textbox(
            Inches(0.6), y,
            Inches(10), Inches(0.4),
        )
        tf = tb_title.text_frame
        tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = (self.title or '').upper()
        r.font.name = 'Inter'; r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['muted'])

        tb_page = slide.shapes.add_textbox(
            self.prs.slide_width - Inches(1.2), y,
            Inches(0.6), Inches(0.4),
        )
        tf = tb_page.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = str(page_num)
        r.font.name = 'Inter'; r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])

    def _title_block(self, slide, title, kicker=None, y_top=Inches(0.8)):
        if kicker:
            tb = slide.shapes.add_textbox(Inches(0.8), y_top, Inches(11.5), Inches(0.3))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = kicker.upper()
            r.font.name = 'Inter'; r.font.size = Pt(12); r.font.bold = True
            r.font.color.rgb = _rgb(self.palette['accent'])
            y_top += Inches(0.45)
        tb = slide.shapes.add_textbox(Inches(0.8), y_top, Inches(11.5), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = 'Inter'; r.font.size = Pt(32); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])
        # thin underline
        under = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), y_top + Inches(0.95),
            Inches(0.7), Inches(0.04),
        )
        under.line.fill.background()
        under.fill.solid()
        under.fill.fore_color.rgb = _rgb(self.palette['accent'])

    def _bullet_block(self, slide, items, x=Inches(0.8), y=Inches(2.2), w=Inches(11.5), h=Inches(4.5)):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14)
            r = p.add_run()
            r.text = '•  ' + item
            r.font.name = 'Inter'; r.font.size = Pt(20)
            r.font.color.rgb = _rgb(self.palette['text'])

    def _notes(self, slide, text):
        if not text: return
        slide.notes_slide.notes_text_frame.text = text

    # ─── Public builders ─────────────────────────────────────────────────

    def cover(self, notes=None):
        slide = self._new_slide(with_footer=False)
        # accent side-bar on the left
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(0.4), self.prs.slide_height,
        )
        accent.line.fill.background()
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(self.palette['accent'])

        # Big title
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(11), Inches(2.2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = self.title
        r.font.name = 'Inter'; r.font.size = Pt(48); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])

        if self.subtitle:
            tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.6), Inches(11), Inches(0.8))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run(); r2.text = self.subtitle
            r2.font.name = 'Inter'; r2.font.size = Pt(22)
            r2.font.color.rgb = _rgb(self.palette['muted'])

        # Metadata block
        meta_lines = [x for x in [self.author, self.institution, self.date] if x]
        if meta_lines:
            tb3 = slide.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(11), Inches(1.2))
            tf3 = tb3.text_frame
            for i, line in enumerate(meta_lines):
                p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
                r = p.add_run(); r.text = line
                r.font.name = 'Inter'; r.font.size = Pt(14)
                r.font.color.rgb = _rgb(self.palette['text'])

        self._notes(slide, notes)
        return slide

    def agenda(self, items, notes=None):
        slide = self._new_slide()
        self._title_block(slide, 'Agenda')
        # numbered list, large
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5))
        tf = tb.text_frame; tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(20)
            n = p.add_run(); n.text = f'{i+1:02d}.  '
            n.font.name = 'Inter'; n.font.size = Pt(26); n.font.bold = True
            n.font.color.rgb = _rgb(self.palette['accent'])
            r = p.add_run(); r.text = item
            r.font.name = 'Inter'; r.font.size = Pt(26)
            r.font.color.rgb = _rgb(self.palette['text'])
        self._notes(slide, notes)
        return slide

    def section(self, title, subtitle=None, notes=None):
        slide = self._new_slide()
        # Full-surface bg.
        surface = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height,
        )
        surface.line.fill.background()
        surface.fill.solid()
        surface.fill.fore_color.rgb = _rgb(self.palette['surface'])
        # Keep footer on top by inserting bg underneath existing shapes.
        spTree = surface._element.getparent()
        spTree.remove(surface._element)
        spTree.insert(2, surface._element)
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.5))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = 'Inter'; r.font.size = Pt(54); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = subtitle
            r2.font.name = 'Inter'; r2.font.size = Pt(22)
            r2.font.color.rgb = _rgb(self.palette['muted'])
        self._notes(slide, notes)
        return slide

    def text_slide(self, title, bullets, kicker=None, notes=None):
        slide = self._new_slide()
        self._title_block(slide, title, kicker=kicker)
        self._bullet_block(slide, bullets)
        self._notes(slide, notes)
        return slide

    def two_column(self, title, *, left_title, left, right_title, right, kicker=None, notes=None):
        slide = self._new_slide()
        self._title_block(slide, title, kicker=kicker)
        # Left
        tb_l = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.7), Inches(0.5))
        p = tb_l.text_frame.paragraphs[0]
        r = p.add_run(); r.text = left_title
        r.font.name = 'Inter'; r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])
        self._bullet_block(slide, left, x=Inches(0.8), y=Inches(2.8), w=Inches(5.7), h=Inches(4))
        # Right
        tb_r = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.7), Inches(0.5))
        p = tb_r.text_frame.paragraphs[0]
        r = p.add_run(); r.text = right_title
        r.font.name = 'Inter'; r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])
        self._bullet_block(slide, right, x=Inches(6.8), y=Inches(2.8), w=Inches(5.7), h=Inches(4))
        self._notes(slide, notes)
        return slide

    def big_stat(self, value, label, caption=None, notes=None):
        slide = self._new_slide()
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(value)
        r.font.name = 'Inter'; r.font.size = Pt(120); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['accent'])
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.2))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label
        r2.font.name = 'Inter'; r2.font.size = Pt(28)
        r2.font.color.rgb = _rgb(self.palette['primary'])
        if caption:
            tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6))
            tf3 = tb3.text_frame
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run(); r3.text = caption
            r3.font.name = 'Inter'; r3.font.size = Pt(14)
            r3.font.color.rgb = _rgb(self.palette['muted'])
        self._notes(slide, notes)
        return slide

    def quote(self, text, attribution=None, notes=None):
        slide = self._new_slide()
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(3))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f'"{text}"'
        r.font.name = 'Inter'; r.font.size = Pt(36); r.font.italic = True
        r.font.color.rgb = _rgb(self.palette['primary'])
        if attribution:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(20)
            r2 = p2.add_run(); r2.text = f'— {attribution}'
            r2.font.name = 'Inter'; r2.font.size = Pt(18)
            r2.font.color.rgb = _rgb(self.palette['muted'])
        self._notes(slide, notes)
        return slide

    def thanks(self, message='Gracias.', contact=None, notes=None):
        slide = self._new_slide(with_footer=False)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(0.4), self.prs.slide_height,
        )
        accent.line.fill.background()
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(self.palette['accent'])
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(11), Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = message
        r.font.name = 'Inter'; r.font.size = Pt(72); r.font.bold = True
        r.font.color.rgb = _rgb(self.palette['primary'])
        if contact:
            tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(11), Inches(0.8))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run(); r2.text = contact
            r2.font.name = 'Inter'; r2.font.size = Pt(16)
            r2.font.color.rgb = _rgb(self.palette['muted'])
        self._notes(slide, notes)
        return slide

    def save(self, out_path):
        self.prs.save(out_path)
